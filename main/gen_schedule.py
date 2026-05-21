# gen_schedule.py
# -*- coding: utf-8 -*-

import io
import re
from dataclasses import dataclass
from typing import Dict, Optional, Tuple

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---- 色（社内スケジュールPPTでよく使われるRGB定義を流用できる形） ----
PALETTE = {
    "Cyan":      RGBColor(0, 169, 224),
    "Gray":      RGBColor(45, 45, 45),
    "Magenta":   RGBColor(218, 24, 132),
    "Turquoise": RGBColor(0, 178, 169),
    "Green":     RGBColor(120, 190, 32),
    "Yellow":    RGBColor(238, 220, 0),
    "Orange":    RGBColor(255, 106, 19),
    "Purple":    RGBColor(128, 49, 167),
    "LightGray": RGBColor(244, 244, 244),
}

STATUS_COLOR = {
    "E": PALETTE["Green"],
    "M": PALETTE["Cyan"],
    "B": PALETTE["Orange"],
}

# 凡例の語彙（テンプレにある想定）
SHAPE_KIND_TO_AUTOSHAPE = {
    "Sample":        MSO_SHAPE.OVAL,
    "Milestone":     MSO_SHAPE.OVAL,
    "Decision point":MSO_SHAPE.DIAMOND,
    "ISO_Events":    MSO_SHAPE.RECTANGLE,
}

@dataclass
class PlotArea:
    left: int
    top: int
    width: int
    height: int

def _parse_ym(s: str) -> Tuple[int, int]:
    """YYYY-MM or YYYY/MM を (year, month) に"""
    s = str(s).strip()
    m = re.match(r"^(\d{4})[-/](\d{1,2})$", s)
    if not m:
        raise ValueError(f"日付形式が不正です: {s}（YYYY-MM または YYYY/MM）")
    y = int(m.group(1))
    mo = int(m.group(2))
    if not (1 <= mo <= 12):
        raise ValueError(f"月が不正です: {s}")
    return y, mo

def _month_index(y: int, m: int, y0: int, m0: int) -> int:
    return (y - y0) * 12 + (m - m0)

def _find_plot_area(slide) -> Optional[PlotArea]:
    """
    テンプレ上に名前 or alt-text で 'PLOT_AREA' を含む図形があれば、それをプロット領域にする
    （最強に自動化するなら、テンプレに透明の矩形を置いて名前を PLOT_AREA にしておくのがオススメ）
    """
    for shp in slide.shapes:
        nm = (shp.name or "").lower()
        if "plot_area" in nm or "plot area" in nm:
            return PlotArea(shp.left, shp.top, shp.width, shp.height)
    return None

def _default_plot_area(prs: Presentation,
                       margin_left=Inches(1.0),
                       margin_top=Inches(1.4),
                       margin_right=Inches(0.8),
                       margin_bottom=Inches(1.0)) -> PlotArea:
    sw = prs.slide_width
    sh = prs.slide_height
    left = int(margin_left)
    top = int(margin_top)
    width = int(sw - margin_left - margin_right)
    height = int(sh - margin_top - margin_bottom)
    return PlotArea(left, top, width, height)

def _rgb_from_any(x) -> Optional[RGBColor]:
    if x is None or (isinstance(x, float) and pd.isna(x)):
        return None
    s = str(x).strip()
    if not s:
        return None
    # "#RRGGBB" or "RRGGBB"
    if re.match(r"^#?[0-9a-fA-F]{6}$", s):
        s = s.replace("#", "")
        return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
    # "R,G,B"
    if re.match(r"^\d{1,3}\s*,\s*\d{1,3}\s*,\s*\d{1,3}$", s):
        r,g,b = [int(v) for v in s.split(",")]
        return RGBColor(r,g,b)
    # "Green" 等のパレット名
    if s in PALETTE:
        return PALETTE[s]
    return None

def _add_text(slide, x, y, w, h, text, font_pt=8, color=PALETTE["Gray"], bold=False):
    tx = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_pt)
    p.font.bold = bool(bold)
    p.font.color.rgb = color
    return tx

def _add_bar(slide, x0, y, x1, h, color: RGBColor):
    w = max(1, int(x1-x0))
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(min(x0,x1)), int(y), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    return shp

def _add_marker(slide, x, y, size, color: RGBColor, kind: str):
    shape = SHAPE_KIND_TO_AUTOSHAPE.get(kind, MSO_SHAPE.OVAL)
    shp = slide.shapes.add_shape(shape, int(x), int(y), int(size), int(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    return shp

def generate_pptx_bytes(
    template_pptx_bytes: bytes,
    df: pd.DataFrame,
    start_ym: str = "2026-01",
    end_ym: str = "2030-12",
    slide_index: int = 0,
    lane_pitch_pt: float = 14.0,
    bar_height_pt: float = 6.0,
    label_font_pt: int = 8,
) -> bytes:
    """
    template pptx(bytes) + df -> out pptx(bytes)
    df 必須列: lane,type,label,start,end,status
    任意列: shape,color,note
    """
    required = ["lane","type","label","start","end","status"]
    for c in required:
        if c not in df.columns:
            raise ValueError(f"DataFrameに必須列 '{c}' がありません: {required}")

    prs = Presentation(io.BytesIO(template_pptx_bytes))
    slide = prs.slides[slide_index]

    pa = _find_plot_area(slide)
    if pa is None:
        pa = _default_plot_area(prs)

    y0,m0 = _parse_ym(start_ym)
    y1,m1 = _parse_ym(end_ym)
    total_months = _month_index(y1,m1,y0,m0) + 1
    if total_months <= 0:
        raise ValueError("start_ym/end_ym の範囲が不正です")

    # lane順を出現順で確定
    lanes = []
    for v in df["lane"].astype(str).tolist():
        if v not in lanes:
            lanes.append(v)
    lane_to_i = {ln:i for i,ln in enumerate(lanes)}

    pitch = int(Pt(lane_pitch_pt))
    bar_h = int(Pt(bar_height_pt))
    marker_size = int(Pt(max(6.0, bar_height_pt*1.8)))

    def x_from_ym(ym: str) -> int:
        yy,mm = _parse_ym(ym)
        k = _month_index(yy,mm,y0,m0)
        k = max(0, min(k, total_months-1))
        denom = (total_months-1) if total_months>1 else 1
        return int(pa.left + pa.width * (k/denom))

    def y_from_lane(lane: str) -> int:
        return int(pa.top + lane_to_i.get(lane,0)*pitch)

    # 描画
    for _, r in df.iterrows():
        lane = str(r["lane"])
        typ = str(r["type"]).strip().lower()   # bar / point
        label = r.get("label","")
        status = str(r.get("status","M")).strip().upper()

        # 色決定：color列があれば優先、無ければstatus色
        color = _rgb_from_any(r.get("color")) or STATUS_COLOR.get(status, PALETTE["Cyan"])

        x0m = x_from_ym(str(r["start"]))
        x1m = x_from_ym(str(r["end"]))
        y = y_from_lane(lane)

        if typ == "bar":
            _add_bar(slide, x0m, y, x1m, bar_h, color)
            _add_text(slide, min(x0m,x1m), y-int(Pt(10)), int(Inches(2.6)), int(Pt(12)),
                      label, font_pt=label_font_pt, color=PALETTE["Gray"])
            note = r.get("note", "")
            if isinstance(note, str) and note.strip():
                _add_text(slide, max(x0m,x1m)+int(Pt(2)), y-int(Pt(10)), int(Inches(3.0)), int(Pt(12)),
                          note, font_pt=max(6,label_font_pt-1), color=PALETTE["Gray"])

        elif typ == "point":
            kind = str(r.get("shape","Milestone")).strip() or "Milestone"
            mx = int(x0m - marker_size//2)
            my = int(y - marker_size//2 + bar_h//2)
            _add_marker(slide, mx, my, marker_size, color, kind)
            _add_text(slide, mx+marker_size+int(Pt(2)), y-int(Pt(10)), int(Inches(3.0)), int(Pt(12)),
                      label, font_pt=label_font_pt, color=PALETTE["Gray"])

        else:
            # 不明typeは無視
            continue

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()